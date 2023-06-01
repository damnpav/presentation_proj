from pptx import Presentation
from pptx.util import Inches

# Create a new Presentation object
prs = Presentation()

# Slide 1: Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "SQL Анализ"
subtitle.text = "Сравнительный анализ двух источников данных"

# Slide 2: Creating Views
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Создание представлений"
content = slide.placeholders[1]
content.text = '''
CREATE VIEW org AS
SELECT id AS org_id, countryresident_name AS org_countryresident_name, inn AS org_inn, kpp AS org_kpp, legalclassification_shortname AS org_legalclassification_shortname, clienttype_oud AS org_clienttype_oud, crm_macroindustry_name AS org_crm_macroindustry_name, okatocode AS org_okatocode, okvedcode_main AS org_okvedcode_main FROM prx_yepk_oud_5_custom_cib_pcapoud.sdp_p4d_ucp_r_organization;
CREATE VIEW const AS
SELECT id AS const_id, rezident_type AS const_rezident_type, inn AS const_inn, kpp AS const_kpp, OPF AS const_OPF, ref_holding AS const_ref_holding, branch AS const_branch, okato AS const_okato, okved_code AS const_okved_code FROM prx_yepk_oud_5_custom_cib_pcapoud.sdp_p4d_ucp_r_org_const;
'''

# Slide 3: Checking Count of IDs
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Проверка количества ID"
content = slide.placeholders[1]
content.text = '''
SELECT COUNT(DISTINCT id) FROM org;
SELECT COUNT(DISTINCT id) FROM const;
'''

# Slide 4: Joining sources
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Объединение источников"
content = slide.placeholders[1]
content.text = '''
CREATE VIEW joined AS
SELECT * FROM org JOIN const ON org.org_id = const.const_id;
'''

# Slide 5: Check size of joined table
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Проверка размера объединенной таблицы"
content = slide.placeholders[1]
content.text = 'SELECT COUNT(*) FROM joined;'

# Slide 6: Checking for different values
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Проверка наличия разных значений"
content = slide.placeholders[1]
content.text = '''
CREATE VIEW buf AS
SELECT LOWER(org_countryresident_name) AS org_countryresident_name,
CASE WHEN const_rezident_type = 'Резидент' THEN 'россия' ELSE const_rezident_type END AS const_rezident_type
FROM joined;
SELECT COUNT(*) FROM buf WHERE org_countryresident_name != const_rezident_type;
SELECT * FROM buf WHERE org_countryresident_name != const_rezident_type LIMIT 5;
SELECT COUNT(*) FROM joined WHERE org_kpp != const_kpp;
SELECT COUNT(*) FROM joined WHERE org_legalclassification_shortname != const_OPF;
SELECT org_legalclassification_shortname, const_OPF FROM joined WHERE org_legalclassification_shortname != const_OPF LIMIT 5;
'''

# Slide 7: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Заключение"
content = slide.placeholders[1]
content.text = "Во всех неправильных случаях |org_legalclassification_shortname| просто не заполнен."

# Save the presentation
prs.save("sql_analysis.pptx")
